#!/usr/bin/env python3
"""로우 바둑이 GDD → PPTX 생성 스크립트"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 팔레트 ──
BG_DARK = RGBColor(0x0D, 0x1B, 0x2A)       # Midnight Navy
BG_GREEN = RGBColor(0x1A, 0x3A, 0x2A)      # Deep Felt Green
GOLD = RGBColor(0xC9, 0xA8, 0x4C)          # Premium Gold
LIGHT_GOLD = RGBColor(0xFF, 0xD7, 0x00)    # Vibrant Gold
WHITE = RGBColor(0xF5, 0xF0, 0xE8)         # Light Cream
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
SOFT_RED = RGBColor(0xCC, 0x33, 0x33)
TEAL = RGBColor(0x2E, 0x8B, 0x8B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE,
                bold=False, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=16, color=WHITE,
                   line_spacing=1.3, font_name="맑은 고딕"):
    """여러 줄 텍스트 추가"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, sz, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz or font_size)
        p.font.color.rgb = clr or color
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(4)
    return txBox

def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    """테이블 추가 - data는 2D 리스트"""
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.name = "맑은 고딕"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                else:
                    paragraph.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x2A)
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    return table_shape

def section_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_GREEN)
    add_textbox(slide, 1, 2.5, 11, 1.5, title, 44, GOLD, True, PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, 1, 4.2, 11, 1, subtitle, 20, WHITE, False, PP_ALIGN.CENTER)
    # 구분선
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4), Inches(5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()
    return slide

# ════════════════════════════════════════
# 슬라이드 1: 타이틀
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 1, 1.5, 11, 1.5, "로우 바둑이 (Low Baduki)", 52, GOLD, True, PP_ALIGN.CENTER)
add_textbox(slide, 1, 3.2, 11, 0.8, "Game Design Document (GDD)", 28, WHITE, False, PP_ALIGN.CENTER)

# 구분선
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.2), Inches(5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = GOLD
line.line.fill.background()

add_multi_text(slide, 2, 4.8, 9, 2.5, [
    ("SIGIL S3  |  Unity 6.3 LTS  |  모바일(Android/iOS) → PC → 웹", False, 16, GRAY),
    ("", False, 10, GRAY),
    ("작성일: 2026-02-27  |  버전 1.0", False, 14, GRAY),
    ("한국 국내 전용 서비스", False, 14, GOLD),
], font_size=16, color=GRAY)

# ════════════════════════════════════════
# 슬라이드 2: 목차
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.5, 11, 0.8, "목차 (Table of Contents)", 32, GOLD, True)

toc_left = [
    "01. 게임 개요 (Game Overview)",
    "02. 핵심 메커닉 (Core Mechanics)",
    "03. UI/UX 플로우",
    "04. 아트 방향 (Art Direction)",
    "05. 오디오 설계 (Audio Design)",
]
toc_right = [
    "06. 경제 설계 (Economy Design)",
    "07. 콘텐츠 설계 (Content Design)",
    "08. 기술 요구사항",
    "09. 관리자 도구 (Admin Tools)",
    "10. 리스크 및 법률/심의",
]

for i, item in enumerate(toc_left):
    add_textbox(slide, 1.5, 1.8 + i * 0.7, 5, 0.6, item, 18, WHITE)

for i, item in enumerate(toc_right):
    add_textbox(slide, 7.5, 1.8 + i * 0.7, 5, 0.6, item, 18, WHITE)

# ════════════════════════════════════════
# 섹션 1: 게임 개요
# ════════════════════════════════════════
section_slide("01. 게임 개요", "Game Overview")

# High Concept
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.5, 11, 0.6, "High Concept & Core Fantasy", 28, GOLD, True)

# 큰 따옴표 하이라이트
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2))
box.fill.solid()
box.fill.fore_color.rgb = BG_GREEN
box.line.color.rgb = GOLD
box.line.width = Pt(2)
add_textbox(slide, 2, 1.6, 9, 1, '"언제 어디서든 3분 만에 한 판 — 한국 전통 로우 바둑이의 모던 리부트.\nAI 대전부터 친구 대전까지."', 20, LIGHT_GOLD, True, PP_ALIGN.CENTER)

add_multi_text(slide, 1, 3.2, 5.5, 3.5, [
    ("핵심 재미 요소", True, 20, GOLD),
    ("", False, 6, WHITE),
    ("  전략적 긴장감 — 3회 드로우 두뇌 싸움", False, 15, WHITE),
    ("  블러핑의 쾌감 — 심리전 승리", False, 15, WHITE),
    ("  성장 체감 — ELO/랭크 시스템", False, 15, WHITE),
    ("  빠른 만족감 — 3-5분 세션", False, 15, WHITE),
])

add_multi_text(slide, 7, 3.2, 5.5, 3.5, [
    ("USP (차별점)", True, 20, GOLD),
    ("", False, 6, WHITE),
    ("1. 로우 바둑이 전문 — 서브 모드가 아닌 메인", False, 15, WHITE),
    ("2. AI 대전 3단계 — 온디바이스 추론, 서버 비용 $0", False, 15, WHITE),
    ("3. 모던 프리미엄 UX — 10년+ 노후 경쟁사 대비", False, 15, WHITE),
])

# 장르 및 타겟
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.5, 11, 0.6, "장르 & 타겟 유저", 28, GOLD, True)

add_table(slide, 0.8, 1.4, 5.5, 2.2, 5, 2, [
    ["구분", "내용"],
    ["메인 장르", "턴제 카드 게임"],
    ["서브장르", "포커 변형 — 로우볼"],
    ["세션 유형", "캐주얼 (3-5분/판)"],
    ["소셜 구조", "멀티플레이어 (2-4인) + AI"],
], col_widths=[2, 3.5])

add_table(slide, 6.8, 1.4, 6, 3.5, 5, 5, [
    ["세그먼트", "연령", "월 지출", "수익 기여", "우선순위"],
    ["전략형 경쟁자", "20-40대 남성", "$15-80", "40-50%", "1순위"],
    ["사교형 애호가", "30-50대", "$5-20", "20-30%", "2순위"],
    ["향수형 레트로", "50대+", "$2-5", "5-10%", "3순위"],
    ["하드코어 챌린저", "20-30대", "$40-100+", "20-30%", "4순위"],
], col_widths=[1.5, 1.2, 1, 1, 1.3])

add_multi_text(slide, 0.8, 4.2, 12, 2.5, [
    ("핵심 동기 (한국 게이머)", True, 16, GOLD),
    ("  사회적 연결 36.1%  |  성장/진행 24.6%  |  경쟁/실력 19.3%  |  시간 효율 18.2%", False, 14, WHITE),
    ("", False, 6, WHITE),
    ("Early Adopter: 한게임/피망 현재 유저 중 UI 불만층 (20-40대 남성)", False, 14, GRAY),
])

# ════════════════════════════════════════
# 섹션 2: 핵심 메커닉
# ════════════════════════════════════════
section_slide("02. 핵심 메커닉", "Core Mechanics")

# 코어 루프
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.5, 11, 0.6, "코어 루프 (Core Loop)", 28, GOLD, True)

add_multi_text(slide, 0.8, 1.4, 6, 4, [
    ("게임 진행 순서", True, 20, GOLD),
    ("", False, 6, WHITE),
    ("Step 1: 블라인드 설정 (SB/BB)", False, 14, WHITE),
    ("Step 2: 초기 핸드 4장 배분", False, 14, WHITE),
    ("Step 3: 1라운드 베팅 (프리드로우)", False, 14, WHITE),
    ("Step 4: 1차 드로우 (아침) — 0-4장 교환", False, 14, WHITE),
    ("Step 5: 2라운드 베팅", False, 14, WHITE),
    ("Step 6: 2차 드로우 (점심)", False, 14, WHITE),
    ("Step 7: 3라운드 베팅", False, 14, WHITE),
    ("Step 8: 3차 드로우 (저녁)", False, 14, WHITE),
    ("Step 9: 4라운드 베팅", False, 14, WHITE),
    ("Step 10: 쇼다운 — 족보 비교, 승자 팟 획득", False, 14, WHITE),
])

add_multi_text(slide, 7, 1.4, 5.5, 2.5, [
    ("3단계 루프 구조", True, 20, GOLD),
    ("", False, 6, WHITE),
    ("[단기] 판 → 승/패 → 코인 + 경험치", False, 14, TEAL),
    ("[중기] 레벨업 → 코스메틱 → ELO 랭크 승급", False, 14, TEAL),
    ("[장기] 시즌 보상 → 배틀패스 → 콜렉션 완성", False, 14, TEAL),
])

add_table(slide, 7, 4.2, 5.5, 2, 5, 3, [
    ["단계", "목표 시간", "설명"],
    ["매칭 대기", "30초 이내", "ELO ±50 → 확장"],
    ["게임 플레이", "3-5분", "4라운드 + 3드로우"],
    ["결과/보상", "20-30초", "쇼다운 + 보상"],
], col_widths=[1.5, 1.3, 2.7])

# 족보
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.5, 11, 0.6, "족보 (Hand Rankings) — 낮을수록 강함", 28, GOLD, True)

add_table(slide, 0.8, 1.4, 6, 2.5, 5, 4, [
    ["급", "족보명", "설명", "예시"],
    ["1급", "메이드 (Badugi)", "4장 무늬+숫자 모두 다름", "A♠2♣3♦4♥"],
    ["2급", "베이스 (Three-Card)", "3장만 다른 무늬/숫자", "A♠2♣3♦3♥"],
    ["3급", "투베이스 (Two-Card)", "2장씩 쌍으로 중복", "A♠2♣A♦2♥"],
    ["4급", "노페어 (One-Card)", "모두 같은 무늬 or 숫자", "A♠2♠3♠4♠"],
], col_widths=[0.6, 1.8, 1.8, 1.8])

add_multi_text(slide, 7.5, 1.4, 5, 3.5, [
    ("메이드 내 순위 (최강 → 최약)", True, 18, GOLD),
    ("", False, 6, WHITE),
    ("1위  골프 (Perfect): A-2-3-4", False, 15, LIGHT_GOLD),
    ("2위  Second: A-2-3-5", False, 15, WHITE),
    ("3위  Third: A-2-4-5", False, 15, WHITE),
    ("...  가장 높은 카드부터 비교", False, 15, GRAY),
    ("최약  K-high: K-Q-J-10", False, 15, SOFT_RED),
])

add_multi_text(slide, 0.8, 4.3, 12, 2.5, [
    ("Luck vs Skill 비율", True, 18, GOLD),
    ("  종합: Luck 40% / Skill 60%  |  초기 핸드: Luck 60%  |  드로우: Skill 60%  |  베팅/블러핑: Skill 80%", False, 14, WHITE),
    ("", False, 6, WHITE),
    ("블라인드: 연습 10/20  |  일반 50/100  |  하이롤러 500/1,000", False, 14, GRAY),
    ("시간 제한: 베팅 20초, 드로우 15초, 경고 5초", False, 14, GRAY),
])

# ════════════════════════════════════════
# 섹션 3: UI/UX 플로우
# ════════════════════════════════════════
section_slide("03. UI/UX 플로우", "UI/UX Flow")

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.5, 11, 0.6, "전체 화면 맵 (10개 화면)", 28, GOLD, True)

add_multi_text(slide, 0.8, 1.3, 5.5, 5.5, [
    ("핵심 플로우", True, 18, GOLD),
    ("", False, 6, WHITE),
    ("1. 스플래시 → 로그인/소셜 인증", False, 14, WHITE),
    ("2. 신규 유저 → 튜토리얼 (5단계, 스킵 불가)", False, 14, WHITE),
    ("3. 메인 로비 (홈)", False, 14, WHITE),
    ("4. 게임 모드 선택:", False, 14, WHITE),
    ("   - AI 대전 (난이도 3단계)", False, 13, TEAL),
    ("   - 랭크 매치 (ELO 기반 매칭)", False, 13, TEAL),
    ("   - 친구 대전 (방 생성 + 링크 초대)", False, 13, TEAL),
    ("   - 일반 매치 (P1)", False, 13, GRAY),
    ("   - 토너먼트 (P2)", False, 13, GRAY),
    ("5. 게임 테이블 → 쇼다운 → 결과/보상", False, 14, WHITE),
    ("", False, 6, WHITE),
    ("하단 탭바: 홈 / 상점 / 콜렉션 / 랭킹 / 프로필", False, 13, GRAY),
])

add_multi_text(slide, 7, 1.3, 5.5, 5.5, [
    ("주요 화면 상세", True, 18, GOLD),
    ("", False, 6, WHITE),
    ("메인 로비", True, 15, WHITE),
    ("  상단바: 코인/젬 잔액, 프로필", False, 13, WHITE),
    ("  중앙: 게임 모드 버튼", False, 13, WHITE),
    ("  시즌 패스 진행도 + 일일 미션 알림", False, 13, WHITE),
    ("", False, 6, WHITE),
    ("게임 테이블", True, 15, WHITE),
    ("  상대 핸드(뒷면) / 팟 정보 / 드로우 단계", False, 13, WHITE),
    ("  내 핸드(앞면) + 족보 자동 계산", False, 13, WHITE),
    ("  Fold/Check/Call/Raise/All-in 버튼", False, 13, WHITE),
    ("  금액 슬라이더 + 타이머 바", False, 13, WHITE),
    ("  이모지 채팅 (텍스트 X)", False, 13, WHITE),
    ("", False, 6, WHITE),
    ("친구 대전: 방 생성 + 딥링크 초대", False, 13, WHITE),
    ("  baduki://room/{roomCode}", False, 12, TEAL),
])

# ════════════════════════════════════════
# 섹션 4: 아트 & 오디오
# ════════════════════════════════════════
section_slide("04-05. 아트 & 오디오", "Art Direction & Audio Design")

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.5, 11, 0.6, "비주얼 스타일 & 에셋", 28, GOLD, True)

add_multi_text(slide, 0.8, 1.3, 5.5, 3, [
    ("모던 프리미엄 (Modern Premium)", True, 20, GOLD),
    ("", False, 6, WHITE),
    ("깔끔한 레이아웃 + 고급 질감 + 부드러운 애니메이션", False, 14, WHITE),
    ("", False, 6, WHITE),
    ("레퍼런스:", True, 14, WHITE),
    ("  레이아웃: Balatro (깔끔한 카드 배치)", False, 13, GRAY),
    ("  질감: 프리미엄 카지노 앱 (펠트, 골드 악센트)", False, 13, GRAY),
    ("  애니메이션: Marvel SNAP (부드러운 카드 플립)", False, 13, GRAY),
])

# 컬러 팔레트 테이블
add_table(slide, 0.8, 4, 5.5, 2.8, 7, 3, [
    ["용도", "색상명", "HEX"],
    ["배경 (테이블)", "Deep Felt Green", "#1A3A2A"],
    ["배경 (로비)", "Midnight Navy", "#0D1B2A"],
    ["주요 액션 버튼", "Premium Gold", "#C9A84C"],
    ["승리 강조", "Vibrant Gold", "#FFD700"],
    ["카드 배경", "Pure White", "#FFFFFF"],
    ["UI 텍스트", "Light Cream", "#F5F0E8"],
], col_widths=[1.8, 1.8, 1.9])

add_multi_text(slide, 7, 1.3, 5.5, 5, [
    ("에셋 가이드 (Unity AI Generators)", True, 18, GOLD),
    ("", False, 6, WHITE),
    ("  카드 기본 덱: 52장 (256x384px)", False, 13, WHITE),
    ("  카드 뒷면: 5종", False, 13, WHITE),
    ("  코스메틱 카드 스킨: 20종+", False, 13, WHITE),
    ("  테이블 배경: 5종 (1920x1080px)", False, 13, WHITE),
    ("  UI 아이콘: 50종+ (64x64px)", False, 13, WHITE),
    ("  아바타: 20종 (128x128px)", False, 13, WHITE),
    ("  이모지: 30종 (64x64px)", False, 13, WHITE),
    ("", False, 6, WHITE),
    ("Sprite Atlas (ASTC 6x6 압축 → 메모리 75% 절감)", False, 13, TEAL),
    ("", False, 10, WHITE),
    ("오디오", True, 18, GOLD),
    ("  BGM: 재즈 피아노(로비), 미니멀 일렉트로닉(게임)", False, 13, WHITE),
    ("  SFX: 카드/칩/버튼 13종 (P0) — Unity AI Generators", False, 13, WHITE),
    ("  기본 음량: BGM 40% / SFX 70%", False, 13, GRAY),
])

# ════════════════════════════════════════
# 섹션 6: 경제 설계
# ════════════════════════════════════════
section_slide("06. 경제 설계", "Economy Design")

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.5, 11, 0.6, "통화 시스템 & 수익화", 28, GOLD, True)

add_table(slide, 0.8, 1.3, 5.5, 2.2, 5, 3, [
    ["통화", "획득 방법", "사용처"],
    ["코인 (기본)", "일일 보너스, 승리, 광고", "테이블 참가비, 코스메틱"],
    ["젬 (프리미엄)", "IAP 구매, 배틀패스", "코스메틱, 탄 구매"],
    ["BP 포인트", "게임 참가, 미션", "배틀패스 보상 해금"],
    ["이벤트 토큰", "이벤트 참여", "이벤트 상점"],
], col_widths=[1.3, 2.2, 2])

add_table(slide, 7, 1.3, 5.5, 2.5, 5, 3, [
    ["수익원", "비중", "설명"],
    ["코스메틱 IAP", "60%", "카드 스킨, 보드, 이모지, 프로필"],
    ["리워드 광고", "20%", "패배 후 코인, 일일 무료"],
    ["배틀패스", "15%", "시즌 프리미엄 트랙"],
    ["VIP 구독", "5%", "광고 제거 + 일일 보너스"],
], col_widths=[1.5, 0.8, 3.2])

add_multi_text(slide, 0.8, 3.8, 12, 3.2, [
    ("탄 시스템 (피처 플래그 의존)", True, 16, SOFT_RED),
    ("  심의 시: FEATURE_TAN_SYSTEM = false (완전 비활성화)  |  서비스 시: 관리자 툴에서 ON", False, 13, WHITE),
    ("", False, 8, WHITE),
    ("인앱 구매 주요 항목", True, 16, GOLD),
    ("  카드 스킨 4종: ₩1,990~₩2,490 / 번들 ₩6,990  |  보드 테마: ₩1,490~₩1,990", False, 13, WHITE),
    ("  젬 충전: ₩1,200(100젬) ~ ₩59,000(7,000젬+40%)  |  배틀패스: ₩4,900/시즌", False, 13, WHITE),
    ("  VIP 구독: ₩2,900/월 (광고제거+일일 100젬)  |  탄 콤보 팩: ₩2,490", False, 13, WHITE),
])

# LTV/ARPU
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.5, 11, 0.6, "LTV/ARPU 목표 & ELO 시스템", 28, GOLD, True)

add_table(slide, 0.8, 1.3, 5.5, 3.5, 9, 3, [
    ["지표", "목표값", "기준"],
    ["D1 리텐션", "40%+", "소프트 런치"],
    ["D7 리텐션", "20%+", "소프트 런치"],
    ["D30 리텐션", "10%+", "정식 출시 1개월"],
    ["ARPU (월)", "$1-3", "정식 출시 3개월"],
    ["ARPPU (월)", "$15-25", "Dolphin 기준"],
    ["유료 전환율", "3%+", "소프트 런치"],
    ["DAU/MAU", "20%+", "정식 출시 3개월"],
    ["BEP", "MAU 10,000", "$20K/월 운영비"],
], col_widths=[1.5, 1.3, 2.7])

add_table(slide, 7, 1.3, 5.5, 2.8, 7, 3, [
    ["ELO 범위", "티어명", "K값"],
    ["900 미만", "브론즈", "64 (신규)"],
    ["900-1100", "실버", "32"],
    ["1100-1300", "골드", "32"],
    ["1300-1500", "플래티넘", "32"],
    ["1500-1700", "다이아몬드", "16"],
    ["1700+", "레전드", "16"],
], col_widths=[1.5, 1.5, 2.5])

add_multi_text(slide, 7, 4.5, 5.5, 2.5, [
    ("시즌 리셋 공식", True, 14, GOLD),
    ("  리셋 ELO = 900 + (이전 ELO - 900) x 0.75", False, 13, TEAL),
    ("  예: 1400 → 1275  |  1700 → 1500", False, 12, GRAY),
])

# ════════════════════════════════════════
# 섹션 7: 콘텐츠 설계
# ════════════════════════════════════════
section_slide("07. 콘텐츠 설계", "Content Design")

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.5, 11, 0.6, "게임 모드 & AI 봇 & 시즌", 28, GOLD, True)

add_table(slide, 0.8, 1.3, 6, 2.8, 7, 4, [
    ["모드", "우선순위", "출시 시기", "ELO"],
    ["AI 대전", "P0", "출시", "X"],
    ["튜토리얼", "P0", "출시", "X"],
    ["랭크 매치", "P0", "출시", "O"],
    ["친구 대전", "P0", "출시", "X"],
    ["일반 매치", "P1", "+1개월", "X"],
    ["토너먼트", "P2", "+3개월", "X"],
], col_widths=[1.5, 1, 1.2, 2.3])

add_table(slide, 7.5, 1.3, 5, 2, 4, 4, [
    ["난이도", "알고리즘", "블러핑", "잠금"],
    ["Easy", "Rule-Based LUT", "1%", "없음"],
    ["Medium", "Heuristic+MCTS", "20-30%", "없음"],
    ["Hard", "CFR (Inference Engine)", "GTO", "+1개월"],
], col_widths=[1, 1.5, 0.8, 1.7])

add_multi_text(slide, 0.8, 4.3, 5.5, 3, [
    ("시즌 시스템 (월간)", True, 16, GOLD),
    ("  1일: 시작 (ELO 소프트 리셋 75%)", False, 13, WHITE),
    ("  1-25일: 랭크 매치 + 배틀패스 미션", False, 13, WHITE),
    ("  26-28일: 종료 예고", False, 13, WHITE),
    ("  말일: 결산 + 보상 지급", False, 13, WHITE),
])

add_multi_text(slide, 7, 3.8, 5.5, 3, [
    ("시즌별 특수 모드 (P1 이후)", True, 16, GOLD),
    ("  S2: 인디언 바둑이 (5장→4장 선택)", False, 13, WHITE),
    ("  S3: 하이로우 바둑이 (팟 50:50)", False, 13, WHITE),
    ("  S4: 스피드 바둑이 (턴 5초)", False, 13, WHITE),
    ("", False, 6, WHITE),
    ("배틀패스: 무료 30단계 + 프리미엄 30단계 (₩4,900)", False, 13, GRAY),
    ("소셜: 친구 시스템 + 이모지 채팅 + 리더보드 3종", False, 13, GRAY),
])

# ════════════════════════════════════════
# 섹션 8: 기술 요구사항
# ════════════════════════════════════════
section_slide("08. 기술 요구사항", "Technical Requirements")

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.5, 11, 0.6, "기술 스택 & 아키텍처", 28, GOLD, True)

add_table(slide, 0.8, 1.3, 5.5, 2.5, 6, 3, [
    ["항목", "선택", "근거"],
    ["게임 엔진", "Unity 6.3 LTS", "크로스플랫폼, Unity AI"],
    ["언어", "C# (.NET 7+)", "Unity 표준"],
    ["렌더러", "2D Renderer (URP)", "카드 게임 최적화"],
    ["UI", "UI Toolkit + TMP", "모던 UI"],
    ["AI 추론", "Unity Inference Engine", "온디바이스 CFR"],
], col_widths=[1.2, 2, 2.3])

add_multi_text(slide, 7, 1.3, 5.5, 3, [
    ("Unity AI 3대 기능 적용", True, 18, GOLD),
    ("", False, 6, WHITE),
    ("AI Assistant", True, 14, TEAL),
    ("  코드 생성, 씬 구성, 디버깅", False, 13, WHITE),
    ("AI Generators", True, 14, TEAL),
    ("  카드 52장, 배경 5종, 이모지 30종, SFX", False, 13, WHITE),
    ("Inference Engine", True, 14, TEAL),
    ("  Hard AI 온디바이스 CFR 추론 (ONNX)", False, 13, WHITE),
])

add_multi_text(slide, 0.8, 4.2, 5.5, 3, [
    ("네트워크 아키텍처", True, 16, GOLD),
    ("  Server-Authoritative (서버 권한)", False, 13, WHITE),
    ("  WebSocket (TCP) — 턴 기반 적합", False, 13, WHITE),
    ("  타임아웃: 연결 끊김 20초 → 자동 폴드", False, 13, WHITE),
    ("", False, 6, WHITE),
    ("Phase 1: Unity Netcode ($0)", False, 13, GRAY),
    ("Phase 2: Photon Fusion ($0, 100 CCU)", False, 13, GRAY),
    ("Phase 3: Photon 업그레이드 ($95/년+)", False, 13, GRAY),
])

add_table(slide, 7, 4.2, 5.5, 2.5, 5, 3, [
    ["성능 목표", "값", "비고"],
    ["프레임레이트", "30 FPS (모바일)", "60 FPS (PC)"],
    ["메모리", "< 200MB", ""],
    ["초기 로딩", "< 3초", "스플래시→로비"],
    ["앱 크기", "< 100MB", "초기 다운로드"],
], col_widths=[1.5, 1.5, 2.5])

# 백엔드
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.5, 11, 0.6, "백엔드 & 제3자 서비스 & 비용", 28, GOLD, True)

add_table(slide, 0.8, 1.3, 5.5, 3, 8, 3, [
    ["서비스", "도구", "Phase"],
    ["인증", "Firebase Auth", "1"],
    ["데이터베이스", "Firebase Firestore", "1"],
    ["실시간 동기화", "Firebase Realtime DB", "1"],
    ["분석", "Firebase Analytics", "1"],
    ["서버리스", "Cloud Functions", "1-2"],
    ["관리자 백엔드", "Node.js + Express", "2"],
    ["관리자 DB", "PostgreSQL + Redis", "2"],
], col_widths=[1.5, 2, 2])

add_table(slide, 7, 1.3, 5.5, 3.5, 8, 3, [
    ["서비스", "Phase 1 비용", "Phase 2 비용"],
    ["Firebase", "$0 (Spark)", "$50-200/월"],
    ["Photon Fusion", "-", "$0 (100 CCU)"],
    ["Unity AI 전체", "$0", "$0"],
    ["Google Play", "$25 (1회)", "-"],
    ["Apple Developer", "$99/년", "$99/년"],
    ["AdMob", "$0", "수익 30% 공유"],
    ["GitHub Actions CI/CD", "$0", "$0"],
], col_widths=[1.5, 1.5, 2.5])

add_multi_text(slide, 0.8, 4.8, 12, 2, [
    ("AI 모델 파이프라인 (Hard AI)", True, 16, GOLD),
    ("  Python CFR 학습 → ONNX 내보내기 → Unity Inference Engine → 런타임 추론 (<10ms NPU / <50ms CPU)", False, 13, WHITE),
    ("  모델 입력: 핸드 인코딩 208차원 + 공개 정보 ≈ 230차원  |  목표 크기: 50MB 이하 (float16)", False, 13, GRAY),
])

# ════════════════════════════════════════
# 섹션 9: 관리자 도구
# ════════════════════════════════════════
section_slide("09. 관리자 도구", "Admin Tools")

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.5, 11, 0.6, "관리자 도구 — 서비스와 동등 레벨", 28, GOLD, True)

add_multi_text(slide, 0.8, 1.3, 5.5, 2, [
    ("기술 스택", True, 16, GOLD),
    ("  React 18 + TypeScript + Tailwind CSS", False, 13, WHITE),
    ("  Node.js + Express + PostgreSQL + Redis", False, 13, WHITE),
    ("  WebSocket + Socket.IO (실시간)", False, 13, WHITE),
    ("  Docker + Nginx (배포)", False, 13, WHITE),
])

add_table(slide, 0.8, 3.3, 5.5, 3, 6, 3, [
    ["P0 기능 (출시 전)", "P1 기능 (+1개월)", ""],
    ["유저 조회/상세", "블라인드 구조 조정", ""],
    ["계정 제재 (경고~밴)", "AI 난이도 파라미터", ""],
    ["코인 지급 (로그)", "리워드 테이블 조정", ""],
    ["실시간 모니터링", "ELO K값 조정", ""],
    ["신고 처리", "피처 플래그 관리", ""],
], col_widths=[2.2, 2, 1.3])

add_multi_text(slide, 7, 1.3, 5.5, 3, [
    ("피처 플래그 (핵심)", True, 18, SOFT_RED),
    ("", False, 6, WHITE),
    ("FEATURE_TAN_SYSTEM    OFF (심의 핵심)", False, 13, WHITE),
    ("FEATURE_TOURNAMENT    OFF", False, 13, WHITE),
    ("FEATURE_HARD_AI       OFF", False, 13, WHITE),
    ("FEATURE_SEASON_SPECIAL OFF", False, 13, WHITE),
    ("FEATURE_MAINTENANCE   OFF", False, 13, WHITE),
    ("FEATURE_NEW_USER_GIFT  ON", False, 13, WHITE),
])

add_table(slide, 7, 4.3, 5.5, 2.2, 5, 2, [
    ["역할", "권한"],
    ["Super Admin", "전체 기능 (플래그 포함)"],
    ["운영자", "유저 관리, 신고, 모니터링"],
    ["CS 담당자", "유저 조회, 신고 (밴 불가)"],
    ["분석가", "읽기 전용"],
], col_widths=[1.5, 4])

# ════════════════════════════════════════
# 섹션 10: 리스크 & 법률
# ════════════════════════════════════════
section_slide("10. 리스크 & 법률/심의", "Risk & Legal/Rating")

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.5, 11, 0.6, "리스크 & 법률/심의", 28, GOLD, True)

add_table(slide, 0.8, 1.3, 5.5, 2.5, 6, 3, [
    ["기술 리스크", "심각도", "완화 전략"],
    ["AI 모델 크기 초과", "중", "float16 양자화, 서버 폴백"],
    ["Firebase 비용 급증", "중", "비용 알림, 캐싱, 예산 한도"],
    ["멀티플레이 동기화 버그", "고", "서버 권한 모델 + QA"],
    ["빌드 파편화", "중", "실기기 테스트"],
    ["Photon CCU 초과", "저", "10K DAU 도달 전 업그레이드"],
], col_widths=[1.8, 0.7, 3])

add_table(slide, 7, 1.3, 5.5, 2, 4, 3, [
    ["시장 리스크", "심각도", "완화 전략"],
    ["한게임/피망 대응", "중", "로우 바둑이 전문 포지셔닝"],
    ["도박 이미지", "중", "스킬 게임, 현금 환전 금지"],
    ["유저 지불 의사 미검증", "중", "소프트 런치 A/B 테스트"],
], col_widths=[1.5, 0.7, 3.3])

add_multi_text(slide, 0.8, 4.2, 5.5, 3, [
    ("GRAC 등급 심의", True, 18, GOLD),
    ("  목표 등급: 12세 이용가 (2군)", False, 14, WHITE),
    ("  포지셔닝: '스킬 게임'", False, 14, WHITE),
    ("  심의 비용: 100-150만원", False, 14, WHITE),
    ("  사전 협의 권장", False, 14, WHITE),
])

add_multi_text(slide, 7, 3.8, 5.5, 3, [
    ("규제 준수 (국내)", True, 18, GOLD),
    ("  확률형 아이템 공시 (2024.03~)", False, 14, WHITE),
    ("  구매 한도: 월 50만원, 일 10만원", False, 14, WHITE),
    ("  PIPA: 14세 미만 부모 동의", False, 14, WHITE),
    ("  사행행위 규제법: 현금 환전 절대 금지", False, 14, SOFT_RED),
])

# ════════════════════════════════════════
# 출시 로드맵 요약
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 0.8, 0.5, 11, 0.6, "출시 로드맵 요약", 28, GOLD, True)

# Phase boxes
phases = [
    ("Phase 1", "프로토타입", "AI 대전(Easy/Medium) + 기본 UI\n+ 튜토리얼 + 친구 대전\n\n검증: 20명 플레이 테스트", BG_GREEN),
    ("Phase 2", "소프트 런치 (한국)", "AI 3단계 + 랭크 매치\n+ 코스메틱 IAP + 배틀패스 v1\n\n조건: GRAC 심의 사전 협의", BG_GREEN),
    ("Phase 3", "정식 출시 (한국)", "GRAC 심의 통과 후 출시\n\n마케팅: ASO + 포커 커뮤니티\n+ 쇼폼 콘텐츠", BG_GREEN),
]

for i, (title, subtitle, desc, bg) in enumerate(phases):
    left = 1 + i * 4
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.5), Inches(3.5), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.color.rgb = GOLD
    box.line.width = Pt(2)

    add_textbox(slide, left + 0.2, 1.7, 3.1, 0.5, title, 22, GOLD, True, PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.2, 2.3, 3.1, 0.5, subtitle, 16, WHITE, True, PP_ALIGN.CENTER)

    # 설명
    txBox = slide.shapes.add_textbox(Inches(left + 0.3), Inches(3), Inches(2.9), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in desc.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.font.name = "맑은 고딕"

# Arrow between phases
for i in range(2):
    left = 4.5 + i * 4
    add_textbox(slide, left, 3.2, 1, 0.5, "→", 36, GOLD, True, PP_ALIGN.CENTER)

add_textbox(slide, 1, 6.5, 11, 0.5, "한국 국내 전용 서비스  |  Trine(개발) 진입은 S4 기획 패키지 승인 후", 14, GRAY, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════
# 마지막 슬라이드
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_textbox(slide, 1, 2, 11, 1.5, "로우 바둑이 (Low Baduki)", 44, GOLD, True, PP_ALIGN.CENTER)
add_textbox(slide, 1, 3.5, 11, 0.8, "Game Design Document — SIGIL S3", 24, WHITE, False, PP_ALIGN.CENTER)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.5), Inches(5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = GOLD
line.line.fill.background()

add_textbox(slide, 1, 5, 11, 0.5, "감사합니다", 28, WHITE, False, PP_ALIGN.CENTER)
add_textbox(slide, 1, 6, 11, 0.5, "2026-02-27  |  Version 1.0  |  한국 국내 전용", 14, GRAY, False, PP_ALIGN.CENTER)

# ── 저장 ──
output_path = "/home/damools/business/02-product/projects/baduki/2026-02-27-s3-gdd.pptx"
prs.save(output_path)
print(f"GDD PPTX 생성 완료: {output_path}")
print(f"슬라이드 수: {len(prs.slides)}")
